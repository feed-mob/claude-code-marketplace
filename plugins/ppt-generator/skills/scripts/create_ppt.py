#!/usr/bin/env python3
"""
PPT Generator Script
Creates PowerPoint presentations using python-pptx library.
Supports creating slides with text, images, themes, and advanced styling.
Enhanced with automatic background selection and logo insertion.
"""

import argparse
import json
import random
import re
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.oxml.xmlchemy import OxmlElement
except ImportError:
    print("Error: python-pptx library is required. Install it with: pip install python-pptx")
    sys.exit(1)


# Asset paths for backgrounds and logos
ASSETS_PATH = Path(__file__).parent.parent / "assets"
BACKGROUNDS_PATH = ASSETS_PATH / "backgrounds"
LOGOS_PATH = ASSETS_PATH / "logos"

# Available background images
BACKGROUND_IMAGES = [
    "bg1.jpg", "bg2.png", "bg3.png", "bg4.jpg", "bg5.jpg"
]

# Available logo images
LOGO_IMAGES = [
    "fm_logo_1.png", "fm_logo_2.png", "fm_logo_3.png"
]

# Content-based background selection mapping
CONTENT_BACKGROUNDS = {
    # Business/Professional content
    "business": ["bg4.jpg", "bg5.jpg"],
    "professional": ["bg4.jpg", "bg5.jpg"],
    "corporate": ["bg4.jpg", "bg5.jpg"],
    "meeting": ["bg5.jpg"],
    "presentation": ["bg4.jpg", "bg5.jpg"],

    # Technology/Digital content
    "technology": ["bg2.png", "bg3.png"],
    "digital": ["bg2.png", "bg3.png"],
    "platform": ["bg2.png", "bg3.png"],
    "software": ["bg2.png", "bg3.png"],
    "app": ["bg2.png", "bg3.png"],
    "online": ["bg2.png", "bg3.png"],

    # Analytics/Data content
    "analytics": ["bg3.png"],
    "data": ["bg3.png"],
    "metrics": ["bg3.png"],
    "performance": ["bg3.png"],
    "reporting": ["bg3.png"],

    # Marketing/Growth content
    "marketing": ["bg1.jpg"],
    "growth": ["bg1.jpg"],
    "campaign": ["bg1.jpg"],
    "audience": ["bg1.jpg"],
    "conversion": ["bg1.jpg"],

    # General/Creative content
    "overview": ["bg1.jpg", "bg4.jpg"],
    "introduction": ["bg1.jpg", "bg4.jpg"],
    "summary": ["bg4.jpg", "bg5.jpg"],
    "conclusion": ["bg4.jpg", "bg5.jpg"],
    "features": ["bg2.png", "bg3.png"],
    "benefits": ["bg1.jpg", "bg4.jpg"],
    "solutions": ["bg2.png", "bg4.jpg"]
}


def select_background_by_content(title: str, content: List[str] = None) -> Optional[str]:
    """Select an appropriate background image based on slide content."""
    title_lower = title.lower()
    content_text = " ".join(content).lower() if content else ""
    full_text = f"{title_lower} {content_text}"

    # Try to match content keywords with background categories
    for category, backgrounds in CONTENT_BACKGROUNDS.items():
        if category in full_text:
            selected_bg = random.choice(backgrounds)
            bg_path = BACKGROUNDS_PATH / selected_bg
            if bg_path.exists():
                return str(bg_path)

    # Fallback to random background
    available_bgs = [bg for bg in BACKGROUND_IMAGES if (BACKGROUNDS_PATH / bg).exists()]
    if available_bgs:
        selected_bg = random.choice(available_bgs)
        return str(BACKGROUNDS_PATH / selected_bg)

    return None


def select_logo_by_content(title: str, content: List[str] = None) -> Optional[str]:
    """Select an appropriate logo image based on slide content."""
    # For now, use random logo selection
    available_logos = [logo for logo in LOGO_IMAGES if (LOGOS_PATH / logo).exists()]
    if available_logos:
        selected_logo = random.choice(available_logos)
        return str(LOGOS_PATH / selected_logo)
    return None


def should_add_background(slide_type: str, title: str, has_existing_background: bool = False) -> bool:
    """Determine if a slide should have a background image."""
    if has_existing_background:
        return False

    # Don't add backgrounds to metrics dashboards or comparison slides
    if slide_type in ['metrics_dashboard', 'comparison']:
        return False

    # Add backgrounds to title slides, section headers, and visual content slides
    if slide_type in ['title', 'section_header', 'visual_content']:
        return True

    # For content slides, add background based on title keywords
    title_lower = title.lower()
    bg_keywords = ['overview', 'introduction', 'summary', 'conclusion', 'features', 'benefits']
    return any(keyword in title_lower for keyword in bg_keywords)


def should_add_logo(slide_type: str, title: str, has_background: bool = False) -> bool:
    """Determine if a slide should have a logo."""
    # Don't add logos to slides that already have backgrounds
    if has_background:
        return False

    # Don't add logos to title slides, section headers, or metrics dashboards
    if slide_type in ['title', 'section_header', 'metrics_dashboard']:
        return False

    # Add logos to content slides, comparison slides, and visual content slides
    return slide_type in ['content', 'comparison', 'visual_content', 'two_column']


def add_logo_to_slide(slide, logo_path: str, position: str = 'bottom_right'):
    """Add a logo to a slide at the specified position."""
    if not Path(logo_path).exists():
        return None

    # Determine slide dimensions (fall back to default 10x7.5 in if unavailable)
    try:
        presentation = slide.part.package.presentation_part.presentation
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
    except AttributeError:
        slide_width = Inches(10)
        slide_height = Inches(7.5)

    logo_width = Inches(1)
    logo_height = Inches(1)
    margin = Inches(0.3)

    max_left = max(margin, slide_width - logo_width - margin)
    max_top = max(margin, slide_height - logo_height - margin)

    # Define logo positions (left, top in EMUs)
    positions = {
        'top_left': (margin, margin),
        'top_right': (max_left, margin),
        'bottom_left': (margin, max_top),
        'bottom_right': (max_left, max_top),
        'center': ((slide_width - logo_width) / 2, (slide_height - logo_height) / 2)
    }

    left, top = positions.get(position, positions['bottom_right'])

    # Add logo with small size (1x1 inch)
    try:
        logo = slide.shapes.add_picture(logo_path, left, top,
                                       logo_width, logo_height)

        # Set transparency for subtle appearance
        try:
            logo.picture_format.transparency = 0.3
        except:
            # Some python-pptx versions may not support this
            pass

        # Send logo to back so it doesn't interfere with content
        logo.z_order = 0
        return logo
    except Exception as e:
        print(f"Warning: Could not add logo {logo_path}: {e}")
        return None


# Design Principles Implementation

def optimize_content_for_6x6(content: List[str]) -> List[str]:
    """Optimize content following the 6x6 rule: max 6 lines per slide, 6 words per line."""
    optimized = []

    for line in content[:6]:  # Max 6 lines per slide
        words = line.split()
        if len(words) > 6:
            # Truncate long lines to 6 words
            truncated = ' '.join(words[:6])
            if len(words) > 6:
                truncated += '...'
            optimized.append(truncated)
        else:
            optimized.append(line)

    return optimized


def enhance_content_hierarchy(content: List[str], color_scheme: str = 'feedmob') -> List[Dict]:
    """Add visual hierarchy to content with emphasis and styling."""
    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])
    enhanced_content = []

    for i, text in enumerate(content):
        # Determine emphasis level based on position and content
        emphasis_level = 'normal'
        text_color = colors['text']
        is_bold = False

        # First item gets emphasis
        if i == 0:
            emphasis_level = 'high'
            text_color = colors['primary']
            is_bold = True

        # Items with action words or metrics get medium emphasis
        action_words = ['improve', 'increase', 'achieve', 'deliver', 'enhance', 'optimize']
        metrics_pattern = re.compile(r'\d+%|\$\d+|\d+\.?\d*[KMB]?')

        if any(word in text.lower() for word in action_words):
            emphasis_level = 'medium'
            text_color = colors['accent1']
            is_bold = True
        elif metrics_pattern.search(text):
            emphasis_level = 'medium'
            text_color = colors['accent2']
            is_bold = True

        enhanced_content.append({
            'text': text,
            'color': text_color,
            'bold': is_bold,
            'emphasis': emphasis_level
        })

    return enhanced_content


def calculate_grid_position(position_type: str, slide_area: str = 'main') -> Tuple[float, float, float, float]:
    """Calculate grid-aligned positions based on design principles."""
    # Grid system: 12 columns x 8 rows
    grid_cols = 12
    grid_rows = 8
    col_width = 10.0 / grid_cols  # 10 inches total width
    row_height = 7.5 / grid_rows  # 7.5 inches total height

    # Define grid positions for different elements
    positions = {
        'title': (1, 0.5, 10, 1.2),  # cols 1-12, rows 1-2
        'subtitle': (1, 1.8, 10, 0.8),  # cols 1-12, rows 3
        'content_left': (1, 2.5, 5.5, 4.5),  # cols 1-7, rows 4-8
        'content_right': (7, 2.5, 4, 4.5),   # cols 8-12, rows 4-8
        'visual_content': (7, 2.5, 4, 4.5),  # cols 8-12, rows 4-8
        'metrics_box': lambda col, row: (col * 3 - 2.5, row * 2 + 0.5, 2.8, 1.8)  # 2x2 grid for metrics
    }

    if position_type == 'metrics_box':
        # Return a function for metrics boxes
        return positions['metrics_box']

    return positions.get(position_type, (1, 1, 8, 5))


def apply_design_principles_to_text(text_frame, content_type: str = 'body', color_scheme: str = 'feedmob'):
    """Apply design principles to text formatting."""
    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            font = run.font

            # Apply typography rules based on design principles
            if content_type == 'title':
                font.size = Pt(44)  # ≥28pt for titles, using 44pt for impact
                font.bold = True
                font.color.rgb = colors['primary']
            elif content_type == 'subtitle':
                font.size = Pt(28)  # ≥28pt for subtitles
                font.color.rgb = colors['text_light']
            elif content_type == 'body':
                font.size = Pt(20)  # ≥18pt for body, using 20pt for readability
                font.color.rgb = colors['text']
            elif content_type == 'emphasis':
                font.size = Pt(22)
                font.bold = True
                font.color.rgb = colors['accent1']

            # Use sans-serif font (Arial) as per design principles
            font.name = 'Arial'


def ensure_white_space(slide_layout: str) -> Dict[str, float]:
    """Calculate proper white space margins based on design principles."""
    # White space: leave at least 20% of slide area as margins
    margins = {
        'left': 0.8,
        'right': 0.8,
        'top': 0.6,
        'bottom': 0.6,
        'content_spacing': 0.3
    }

    # Adjust margins based on slide type
    if slide_layout == 'title':
        margins['top'] = 1.0
        margins['bottom'] = 1.2
    elif slide_layout == 'content':
        margins['content_spacing'] = 0.4
    elif slide_layout == 'visual':
        margins['left'] = 0.6
        margins['right'] = 0.6

    return margins


def apply_60_30_10_color_rule(slide_type: str, color_scheme: str = 'feedmob') -> Dict[str, RGBColor]:
    """Apply the 60-30-10 color rule: 60% primary, 30% secondary, 10% accent."""
    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])

    # Define color usage based on 60-30-10 rule
    color_distribution = {
        'primary': colors['primary'],      # 60% - Main color for backgrounds, titles
        'secondary': colors['secondary'],  # 30% - Secondary color for content areas
        'accent': colors['accent1'],      # 10% - Accent color for highlights and emphasis
        'text': colors['text'],           # Text color
        'background': colors['background'],  # Background color
        'emphasis_high': colors['accent2'],  # Additional accent for high emphasis
        'emphasis_medium': colors['accent3'], # Additional accent for medium emphasis
    }

    # Adjust based on slide type
    if slide_type == 'title':
        # Title slides: more emphasis on primary color
        color_distribution['background'] = colors['primary']
        color_distribution['primary'] = colors['background']
        color_distribution['accent'] = colors['accent2']
    elif slide_type == 'metrics':
        # Metrics dashboards: balanced use of colors for data visualization
        color_distribution['data_primary'] = colors['primary']
        color_distribution['data_secondary'] = colors['accent1']
        color_distribution['data_accent'] = colors['accent2']
    elif slide_type == 'visual':
        # Visual content: emphasis on content clarity
        color_distribution['content_bg'] = colors['background_alt']
        color_distribution['visual_accent'] = colors['accent1']

    return color_distribution


# Color schemes extracted from professional presentations
COLOR_SCHEMES = {
  'feedmob': {
      'primary': RGBColor(5, 141, 199),      # Primary Blue (#058DC7) - extracted from FeedMob theme
      'secondary': RGBColor(21, 129, 88),    # Medium Dark (#158158) - extracted from FeedMob theme
      'accent1': RGBColor(80, 180, 50),     # Secondary Green (#50B432) - extracted from FeedMob theme
      'accent2': RGBColor(237, 86, 27),     # Tertiary Orange (#ED561B) - extracted from FeedMob theme
      'accent3': RGBColor(237, 239, 0),     # Yellow (#EDEF00) - extracted from FeedMob theme
      'accent4': RGBColor(36, 203, 229),    # Cyan (#24CBE5) - extracted from FeedMob theme
      'accent5': RGBColor(100, 229, 114),   # Light Green (#64E572) - extracted from FeedMob theme
      'text': RGBColor(0, 0, 0),           # Black (#000000) - extracted from FeedMob theme
      'background': RGBColor(255, 255, 255), # White (#FFFFFF) - extracted from FeedMob theme
      'background_alt': RGBColor(243, 243, 243), # Light Background Alt (#F3F3F3) - extracted from FeedMob theme
      'text_light': RGBColor(21, 129, 88)   # Medium Dark for secondary text
  }
}


def apply_color_scheme(prs: Presentation, scheme_name: str = 'feedmob'):
    """Apply a color scheme to the presentation."""
    if scheme_name not in COLOR_SCHEMES:
        scheme_name = 'feedmob'

    colors = COLOR_SCHEMES[scheme_name]

    # Store colors for later use in text and shapes
    prs._color_scheme = colors
    return prs


def set_font_style(text_frame, font_name: str = 'Arial', font_size: int = 18,
                  color: Optional[RGBColor] = None, bold: bool = False,
                  italic: bool = False):
    """Apply font styling to text."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            font = run.font
            font.name = font_name
            font.size = Pt(font_size)
            font.bold = bold
            font.italic = italic

            if color:
                font.color.rgb = color


def add_shape_with_text(slide, text: str, left: float, top: float, width: float,
                       height: float, shape_type: str = 'rectangle',
                       fill_color: Optional[RGBColor] = None,
                       text_color: Optional[RGBColor] = None,
                       font_size: int = 18):
    """Add a shape with text to a slide."""
    left_inches = Inches(left)
    top_inches = Inches(top)
    width_inches = Inches(width)
    height_inches = Inches(height)

    if shape_type == 'rectangle':
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_inches, top_inches,
                                      width_inches, height_inches)
    elif shape_type == 'rounded_rectangle':
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_inches, top_inches,
                                      width_inches, height_inches)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_inches, top_inches,
                                      width_inches, height_inches)

    # Set fill color
    if fill_color:
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = fill_color

    # Set text
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)

    # Set text properties
    default_color = text_color or RGBColor(0, 0, 0)
    set_font_style(text_frame, font_size=font_size, color=default_color)

    return shape


def create_gradient_background(slide, start_color: RGBColor, end_color: RGBColor):
    """Add a gradient background to a slide."""
    background = slide.background
    fill = background.fill
    fill.gradient()

    # Set gradient colors
    fill.gradient_stops[0].color.rgb = start_color
    fill.gradient_stops[1].color.rgb = end_color


def add_image_with_effects(slide, image_path: str, left: float, top: float,
                          width: Optional[float] = None, height: Optional[float] = None,
                          shadow: bool = True):
    """Add an image to a slide with optional effects."""
    if not Path(image_path).exists():
        print(f"Warning: Image file not found: {image_path}")
        return None

    left_inches = Inches(left)
    top_inches = Inches(top)

    if width and height:
        width_inches = Inches(width)
        height_inches = Inches(height)
        picture = slide.shapes.add_picture(image_path, left_inches, top_inches,
                                         width=width_inches, height=height_inches)
    else:
        picture = slide.shapes.add_picture(image_path, left_inches, top_inches)

    # Add shadow effect if requested
    if shadow:
        picture.shadow.inherit = False
        picture.shadow.visible = True
        picture.shadow.blur_radius = Inches(0.01)
        picture.shadow.distance = Inches(0.005)
        try:
            picture.shadow.color.rgb = RGBColor(128, 128, 128)
        except AttributeError:
            # Some versions of python-pptx don't support shadow.color
            pass

    return picture


def create_two_column_slide(prs: Presentation, title: str, left_content: List[str],
                           right_content: List[str], color_scheme: str = 'professional'):
    """Create a slide with two columns of content."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])

    # Add left column
    left_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
    left_tf = left_textbox.text_frame
    left_tf.text = left_content[0] if left_content else ""

    for text in left_content[1:]:
        p = left_tf.add_paragraph()
        p.text = text
        p.level = 0

    set_font_style(left_tf, font_size=14, color=colors['text'])

    # Add right column
    right_textbox = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(5))
    right_tf = right_textbox.text_frame
    right_tf.text = right_content[0] if right_content else ""

    for text in right_content[1:]:
        p = right_tf.add_paragraph()
        p.text = text
        p.level = 0

    set_font_style(right_tf, font_size=14, color=colors['text'])

    return slide


def create_title_slide(prs: Presentation, title: str, subtitle: Optional[str] = None,
                      color_scheme: str = 'feedmob', background_image: Optional[str] = None):
    """Create a title slide with enhanced styling following design principles."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])

    # Add background image if provided
    if background_image and Path(background_image).exists():
        add_image_with_effects(slide, background_image, 0, 0, 10, 7.5, shadow=False)

    title_shape = slide.shapes.title
    title_shape.text = title

    # Apply design principles to title
    if hasattr(prs, '_color_scheme'):
        apply_design_principles_to_text(title_shape.text_frame, 'title', color_scheme)

    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle

        # Apply design principles to subtitle
        if hasattr(prs, '_color_scheme'):
            apply_design_principles_to_text(subtitle_shape.text_frame, 'subtitle', color_scheme)

    return slide


def create_content_slide(prs: Presentation, title: str, content: List[str],
                        color_scheme: str = 'feedmob'):
    """Create a content slide with title and bullet points following design principles."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])

    # Apply 6x6 rule to content
    optimized_content = optimize_content_for_6x6(content)
    enhanced_content = enhance_content_hierarchy(optimized_content, color_scheme)

    title_shape = slide.shapes.title
    title_shape.text = title

    # Apply design principles to title
    if hasattr(prs, '_color_scheme'):
        apply_design_principles_to_text(title_shape.text_frame, 'title', color_scheme)

    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame

    if enhanced_content:
        # Add first item with enhanced styling
        first_item = enhanced_content[0]
        tf.text = first_item['text']

        # Apply styling to first paragraph
        paragraph = tf.paragraphs[0]
        for run in paragraph.runs:
            run.font.bold = first_item['bold']
            run.font.color.rgb = first_item['color']

        # Add remaining content with hierarchy
        for item in enhanced_content[1:]:
            p = tf.add_paragraph()
            p.text = item['text']
            p.level = 0

            # Apply individual styling based on hierarchy
            for run in p.runs:
                run.font.bold = item['bold']
                run.font.color.rgb = item['color']
                if item['emphasis'] == 'high':
                    run.font.size = Pt(22)
                elif item['emphasis'] == 'medium':
                    run.font.size = Pt(21)
                else:
                    run.font.size = Pt(20)

    # Apply general body styling
    if hasattr(prs, '_color_scheme'):
        apply_design_principles_to_text(tf, 'body', color_scheme)

    return slide


def create_content_slide_with_background(prs: Presentation, title: str, content: List[str],
                                        color_scheme: str = 'feedmob', background_image: Optional[str] = None):
    """Create a content slide with background image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])

    # Add background image
    if background_image and Path(background_image).exists():
        bg_img = slide.shapes.add_picture(background_image, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_img.z_order = 0

        # Add semi-transparent overlay for better text readability
        overlay = add_shape_with_text(slide, "", 0, 0, 10, 7.5, 'rectangle',
                                      fill_color=RGBColor(255, 255, 255))
        overlay.fill.transparency = 0.2

    # Add title with background
    title_bg = add_shape_with_text(slide, "", 0, 0.3, 10, 1, 'rectangle',
                                   fill_color=colors['primary'])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title_tf = title_box.text_frame
    title_tf.text = title
    set_font_style(title_tf, font_name='Arial', font_size=36,
                  color=RGBColor(255, 255, 255), bold=True)

    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    content_tf = content_box.text_frame
    content_tf.text = content[0] if content else ""

    for text in content[1:]:
        p = content_tf.add_paragraph()
        p.text = text
        p.level = 0

    set_font_style(content_tf, font_name='Arial', font_size=18, color=colors['text'])

    return slide


def create_section_header_slide(prs: Presentation, title: str, subtitle: Optional[str] = None,
                               color_scheme: str = 'professional', background_image: Optional[str] = None):
    """Create a section header slide with enhanced styling."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])

    # Add background image or gradient
    if background_image and Path(background_image).exists():
        bg_img = slide.shapes.add_picture(background_image, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_img.z_order = 0

        # Add semi-transparent overlay
        overlay = add_shape_with_text(slide, "", 0, 0, 10, 7.5, 'rectangle',
                                      fill_color=colors['primary'])
        overlay.fill.transparency = 0.3
    else:
        # Add gradient background
        create_gradient_background(slide, colors['primary'], colors['secondary'])

    # Add title shape
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    title_tf = title_box.text_frame
    title_tf.text = title
    set_font_style(title_tf, font_name='Arial', font_size=48, color=colors['background'], bold=True)

    # Add subtitle if provided
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = subtitle
        set_font_style(subtitle_tf, font_name='Arial', font_size=24, color=colors['background'])

    return slide


def create_comparison_slide(prs: Presentation, title: str, left_title: str, right_title: str,
                           left_points: List[str], right_points: List[str],
                           color_scheme: str = 'professional'):
    """Create a comparison slide with two columns."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])

    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = title
    set_font_style(title_tf, font_name='Arial', font_size=36, color=colors['primary'], bold=True)

    # Add left column with background
    left_bg = add_shape_with_text(slide, "", 0.5, 1.5, 4, 5, 'rounded_rectangle',
                                 fill_color=colors['secondary'])
    left_title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(3.6), Inches(0.8))
    left_title_tf = left_title_box.text_frame
    left_title_tf.text = left_title
    set_font_style(left_title_tf, font_name='Arial', font_size=22, color=colors['primary'], bold=True)

    left_content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(3.6), Inches(3.5))
    left_content_tf = left_content_box.text_frame
    left_content_tf.text = left_points[0] if left_points else ""
    for text in left_points[1:]:
        p = left_content_tf.add_paragraph()
        p.text = text
        p.level = 0
    set_font_style(left_content_tf, font_name='Arial', font_size=16, color=colors['text'])

    # Add right column with background
    right_bg = add_shape_with_text(slide, "", 5, 1.5, 4, 5, 'rounded_rectangle',
                                  fill_color=colors['secondary'])
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(3.6), Inches(0.8))
    right_title_tf = right_title_box.text_frame
    right_title_tf.text = right_title
    set_font_style(right_title_tf, font_name='Arial', font_size=22, color=colors['primary'], bold=True)

    right_content_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.5), Inches(3.6), Inches(3.5))
    right_content_tf = right_content_box.text_frame
    right_content_tf.text = right_points[0] if right_points else ""
    for text in right_points[1:]:
        p = right_content_tf.add_paragraph()
        p.text = text
        p.level = 0
    set_font_style(right_content_tf, font_name='Arial', font_size=16, color=colors['text'])

    return slide


def create_professional_title_slide(prs: Presentation, title: str, subtitle: Optional[str] = None,
                                   color_scheme: str = 'feedmob', background_image: Optional[str] = None):
    """Create a professional title slide matching company style."""
    slide_layout = prs.slide_layouts[6]  # Blank layout for full control
    slide = prs.slides.add_slide(slide_layout)

    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])

    # Add solid color background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['primary']

    # Add background image if provided
    if background_image and Path(background_image).exists():
        bg_img = slide.shapes.add_picture(background_image, Inches(0), Inches(0), Inches(10), Inches(7.5))
        # Send image to back
        bg_img.z_order = 0

    # Add subtitle if provided (smaller text at bottom)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = subtitle
        set_font_style(subtitle_tf, font_name='Arial', font_size=24,
                      color=colors['background'], bold=False)

    # Add main title (larger, prominent)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2.5))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.word_wrap = True
    set_font_style(title_tf, font_name='Arial', font_size=54,
                  color=colors['background'], bold=True)

    return slide


def create_feedmob_title_slide(prs: Presentation, title: str, subtitle: Optional[str] = None,
                              background_image: Optional[str] = None):
    """Create a FeedMob-branded title slide with enhanced styling."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    colors = COLOR_SCHEMES['feedmob']

    # Add background image or solid color
    if background_image and Path(background_image).exists():
        bg_img = slide.shapes.add_picture(background_image, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_img.z_order = 0

        # Add semi-transparent overlay for better text readability
        overlay = add_shape_with_text(slide, "", 0, 0, 10, 7.5, 'rectangle',
                                      fill_color=RGBColor(5, 141, 199))
        overlay.fill.transparency = 0.3
    else:
        # Gradient background using FeedMob colors
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 90
        fill.gradient_stops[0].color.rgb = colors['primary']
        fill.gradient_stops[1].color.rgb = colors['secondary']

    # Add accent bar at top
    accent_bar = add_shape_with_text(slide, "", 0, 0, 10, 0.2, 'rectangle',
                                     fill_color=colors['accent2'])

    # Add main title with enhanced styling
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.word_wrap = True
    set_font_style(title_tf, font_name='Arial', font_size=52,
                  color=RGBColor(255, 255, 255), bold=True)

    # Add subtitle if provided
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = subtitle
        set_font_style(subtitle_tf, font_name='Arial', font_size=28,
                      color=RGBColor(255, 255, 255), bold=False)

    return slide


def create_feedmob_content_slide(prs: Presentation, title: str, content: List[str],
                                image_path: Optional[str] = None):
    """Create a FeedMob-branded content slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    colors = COLOR_SCHEMES['feedmob']

    # Add white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['background']

    # Add top accent bar
    accent_bar = add_shape_with_text(slide, "", 0, 0, 10, 0.15, 'rectangle',
                                     fill_color=colors['primary'])

    # Add title with background
    title_bg = add_shape_with_text(slide, "", 0, 0.15, 10, 0.9, 'rectangle',
                                   fill_color=colors['background_alt'])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_tf = title_box.text_frame
    title_tf.text = title
    set_font_style(title_tf, font_name='Arial', font_size=36,
                  color=colors['primary'], bold=True)

    if image_path and Path(image_path).exists():
        # Layout with image on right
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5), Inches(5))
        add_image_with_effects(slide, image_path, 6, 1.8, 3.5, 4, shadow=True)
    else:
        # Full width content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))

    # Add content with FeedMob styling
    content_tf = content_box.text_frame
    content_tf.text = content[0] if content else ""

    for i, point in enumerate(content[1:], 1):
        p = content_tf.add_paragraph()
        p.text = point
        p.level = 0

        # Add custom bullet point with FeedMob accent color
        try:
            from pptx.dml.color import RGBColor
            run = p.runs[0] if p.runs else p.add_run()
            # Add bullet point symbol with accent color
            p.text = f"• {point}"
            bullet_run = p.runs[0]
            bullet_run.font.color.rgb = colors['accent1']
        except:
            pass

    set_font_style(content_tf, font_name='Arial', font_size=18, color=colors['text'])

    return slide


def create_visual_content_slide(prs: Presentation, title: str, content_points: List[str],
                                image_path: Optional[str] = None, color_scheme: str = 'feedmob',
                                background_image: Optional[str] = None):
    """Create a visual content slide with image and bullet points following design principles."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])

    # Add background image if provided
    if background_image and Path(background_image).exists():
        bg_img = slide.shapes.add_picture(background_image, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_img.z_order = 0

        # Add semi-transparent overlay for better text readability
        overlay = add_shape_with_text(slide, "", 0, 0, 10, 7.5, 'rectangle',
                                      fill_color=RGBColor(255, 255, 255))
        overlay.fill.transparency = 0.15

    # Get grid-aligned positions
    margins = ensure_white_space('visual')
    title_pos = calculate_grid_position('title')
    content_left_pos = calculate_grid_position('content_left')
    visual_pos = calculate_grid_position('visual_content')

    # Add title with colored bar using grid positioning
    title_bar = add_shape_with_text(slide, "", title_pos[0], title_pos[1] + 0.2,
                                    title_pos[2], 0.15, 'rectangle',
                                    fill_color=colors['accent1'])

    title_box = slide.shapes.add_textbox(Inches(title_pos[0]), Inches(title_pos[1]),
                                        Inches(title_pos[2]), Inches(title_pos[3]))
    title_tf = title_box.text_frame
    title_tf.text = title
    apply_design_principles_to_text(title_tf, 'title', color_scheme)

    # Apply 6x6 rule and visual hierarchy to content
    optimized_content = optimize_content_for_6x6(content_points)
    enhanced_content = enhance_content_hierarchy(optimized_content, color_scheme)

    if image_path and Path(image_path).exists():
        # Add image on the right side using grid positioning
        add_image_with_effects(slide, image_path, visual_pos[0], visual_pos[1],
                             visual_pos[2], visual_pos[3], shadow=True)

        # Add content points on the left using grid positioning
        content_box = slide.shapes.add_textbox(Inches(content_left_pos[0]), Inches(content_left_pos[1]),
                                            Inches(content_left_pos[2]), Inches(content_left_pos[3]))
        content_tf = content_box.text_frame

        if enhanced_content:
            # Add content with visual hierarchy
            for i, item in enumerate(enhanced_content):
                if i == 0:
                    content_tf.text = item['text']
                    paragraph = content_tf.paragraphs[0]
                else:
                    paragraph = content_tf.add_paragraph()
                    paragraph.text = item['text']

                # Apply individual styling
                for run in paragraph.runs:
                    run.font.bold = item['bold']
                    run.font.color.rgb = item['color']
                    if item['emphasis'] == 'high':
                        run.font.size = Pt(22)
                    elif item['emphasis'] == 'medium':
                        run.font.size = Pt(21)
                    else:
                        run.font.size = Pt(20)
                    run.font.name = 'Arial'

                # Add proper spacing between points
                if i > 0:
                    paragraph.space_before = Pt(6)

        # Apply general styling with proper margins
        apply_design_principles_to_text(content_tf, 'body', color_scheme)

    else:
        # Full width content using grid positioning
        content_box = slide.shapes.add_textbox(Inches(content_left_pos[0]), Inches(content_left_pos[1]),
                                            Inches(content_left_pos[2] + visual_pos[2]), Inches(content_left_pos[3]))
        content_tf = content_box.text_frame

        if enhanced_content:
            # Add content with visual hierarchy
            for i, item in enumerate(enhanced_content):
                if i == 0:
                    content_tf.text = item['text']
                    paragraph = content_tf.paragraphs[0]
                else:
                    paragraph = content_tf.add_paragraph()
                    paragraph.text = item['text']

                # Apply individual styling
                for run in paragraph.runs:
                    run.font.bold = item['bold']
                    run.font.color.rgb = item['color']
                    if item['emphasis'] == 'high':
                        run.font.size = Pt(22)
                    elif item['emphasis'] == 'medium':
                        run.font.size = Pt(21)
                    else:
                        run.font.size = Pt(20)
                    run.font.name = 'Arial'

                # Add proper spacing
                if i > 0:
                    paragraph.space_before = Pt(8)

        # Apply general styling
        apply_design_principles_to_text(content_tf, 'body', color_scheme)

    return slide


def create_metrics_dashboard_slide(prs: Presentation, title: str, metrics: List[Dict],
                                  color_scheme: str = 'feedmob'):
    """Create a dashboard-style slide with key metrics following design principles."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES['feedmob'])

    # Get grid-aligned positions
    margins = ensure_white_space('metrics')
    title_pos = calculate_grid_position('title')

    # Add title with proper spacing
    title_box = slide.shapes.add_textbox(Inches(title_pos[0]), Inches(title_pos[1]),
                                        Inches(title_pos[2]), Inches(title_pos[3]))
    title_tf = title_box.text_frame
    title_tf.text = title
    apply_design_principles_to_text(title_tf, 'title', color_scheme)

    # Apply 60-30-10 color rule for metrics
    primary_color = colors['primary']      # 60%
    secondary_color = colors['accent1']    # 30%
    accent_color = colors['accent2']       # 10%

    # Create metric boxes using grid system
    metric_colors = [primary_color, secondary_color, accent_color, colors['accent3']]

    for i, metric in enumerate(metrics[:4]):  # Limit to 4 metrics for clarity
        row = i // 2
        col = i % 2

        # Use grid positioning
        metrics_pos_func = calculate_grid_position('metrics_box')
        left, top, width, height = metrics_pos_func(col + 1, row + 1)

        # Background box with subtle design
        metric_bg = add_shape_with_text(slide, "", left, top, width, height, 'rounded_rectangle',
                                       fill_color=colors['background_alt'])

        # Add subtle border for definition
        metric_bg.line.color.rgb = colors['secondary']
        metric_bg.line.width = Inches(0.01)

        # Metric value with emphasis
        value_text = str(metric.get('value', ''))
        value_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15),
                                            Inches(width - 0.3), Inches(height * 0.6))
        value_tf = value_box.text_frame
        value_tf.text = value_text

        # Apply emphasis styling to metrics
        for paragraph in value_tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.color.rgb = metric_colors[i % len(metric_colors)]
                run.font.name = 'Arial'

        # Metric label with secondary emphasis
        label_text = metric.get('label', '')
        label_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + height * 0.65),
                                            Inches(width - 0.3), Inches(height * 0.3))
        label_tf = label_box.text_frame
        label_tf.text = label_text

        # Apply label styling
        for paragraph in label_tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = colors['text_light']
                run.font.name = 'Arial'

    return slide


def create_blank_slide(prs: Presentation):
    """Create a blank slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    return slide


def add_text_to_slide(slide, text: str, left: float = 1, top: float = 1,
                      width: float = 8, height: float = 1, font_size: int = 18):
    """Add text box to a slide."""
    left_inches = Inches(left)
    top_inches = Inches(top)
    width_inches = Inches(width)
    height_inches = Inches(height)

    textbox = slide.shapes.add_textbox(left_inches, top_inches, width_inches, height_inches)
    text_frame = textbox.text_frame
    text_frame.text = text

    # Set font size
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(font_size)

    return textbox


def add_image_to_slide(slide, image_path: str, left: float = 1, top: float = 1,
                       width: Optional[float] = None, height: Optional[float] = None):
    """Add an image to a slide."""
    if not Path(image_path).exists():
        print(f"Warning: Image file not found: {image_path}")
        return None

    left_inches = Inches(left)
    top_inches = Inches(top)

    if width and height:
        width_inches = Inches(width)
        height_inches = Inches(height)
        return slide.shapes.add_picture(image_path, left_inches, top_inches,
                                       width=width_inches, height=height_inches)
    else:
        return slide.shapes.add_picture(image_path, left_inches, top_inches)


def create_presentation_from_json(json_data: Dict, output_path: str):
    """Create a presentation from JSON data structure with automatic backgrounds and logos."""
    prs = Presentation()

    # Apply color scheme if specified
    color_scheme = json_data.get('color_scheme', 'feedmob')
    apply_color_scheme(prs, color_scheme)

    # Check if auto backgrounds and logos are enabled (default: True)
    auto_backgrounds = json_data.get('auto_backgrounds', True)
    auto_logos = json_data.get('auto_logos', True)

    slides_data = json_data.get('slides', [])

    for slide_data in slides_data:
        slide_type = slide_data.get('type', 'content')
        slide_color_scheme = slide_data.get('color_scheme', color_scheme)

        if slide_type == 'title':
            title = slide_data.get('title', '')
            subtitle = slide_data.get('subtitle')
            background_image = slide_data.get('background_image')
            use_professional = slide_data.get('professional_style', False)

            # Auto-select background if none provided and auto_backgrounds is enabled
            if not background_image and auto_backgrounds:
                background_image = select_background_by_content(title)

            if slide_color_scheme == 'feedmob':
                slide = create_feedmob_title_slide(prs, title, subtitle, background_image)
            elif use_professional or slide_color_scheme == 'feedmob':
                slide = create_professional_title_slide(prs, title, subtitle, slide_color_scheme, background_image)
            else:
                slide = create_title_slide(prs, title, subtitle, slide_color_scheme, background_image)

        elif slide_type == 'visual_content':
            title = slide_data.get('title', '')
            content_points = slide_data.get('content', [])
            image_path = slide_data.get('image')
            if isinstance(content_points, str):
                content_points = [content_points]

            # Auto-select background if auto_backgrounds is enabled
            background_image = None
            if auto_backgrounds and should_add_background(slide_type, title):
                background_image = select_background_by_content(title, content_points)

            slide = create_visual_content_slide(prs, title, content_points, image_path, slide_color_scheme, background_image)

            # Add logo if auto_logos is enabled and no background was added
            if auto_logos and should_add_logo(slide_type, title, bool(background_image)):
                logo_path = select_logo_by_content(title, content_points)
                if logo_path:
                    add_logo_to_slide(slide, logo_path, 'bottom_right')

        elif slide_type == 'metrics_dashboard':
            title = slide_data.get('title', '')
            metrics = slide_data.get('metrics', [])
            slide = create_metrics_dashboard_slide(prs, title, metrics, slide_color_scheme)

            # Add logo if auto_logos is enabled (metrics dashboards don't have backgrounds)
            if auto_logos and should_add_logo(slide_type, title):
                logo_path = select_logo_by_content(title)
                if logo_path:
                    add_logo_to_slide(slide, logo_path, 'bottom_right')

        elif slide_type == 'content':
            title = slide_data.get('title', '')
            content = slide_data.get('content', [])
            if isinstance(content, str):
                content = [content]

            # Auto-select background if auto_backgrounds is enabled
            background_image = None
            if auto_backgrounds and should_add_background(slide_type, title):
                background_image = select_background_by_content(title, content)

            if slide_color_scheme == 'feedmob' and 'image' in slide_data:
                # Use FeedMob content slide with image
                image_path = slide_data['image']
                slide = create_feedmob_content_slide(prs, title, content, image_path)
            elif background_image:
                # Create content slide with background
                slide = create_content_slide_with_background(prs, title, content, slide_color_scheme, background_image)
            else:
                slide = create_content_slide(prs, title, content, slide_color_scheme)

                # Add image if specified
                if 'image' in slide_data:
                    image_path = slide_data['image']
                    add_image_to_slide(slide, image_path,
                                     left=slide_data.get('image_left', 5),
                                     top=slide_data.get('image_top', 2))

            # Add logo if auto_logos is enabled and no background was added
            if auto_logos and should_add_logo(slide_type, title, bool(background_image)):
                logo_path = select_logo_by_content(title, content)
                if logo_path:
                    add_logo_to_slide(slide, logo_path, 'bottom_right')

        elif slide_type == 'section_header':
            title = slide_data.get('title', '')
            subtitle = slide_data.get('subtitle')

            # Auto-select background if auto_backgrounds is enabled
            background_image = None
            if auto_backgrounds and should_add_background(slide_type, title):
                background_image = select_background_by_content(title)

            slide = create_section_header_slide(prs, title, subtitle, slide_color_scheme, background_image)

        elif slide_type == 'comparison':
            title = slide_data.get('title', '')
            left_title = slide_data.get('left_title', '')
            right_title = slide_data.get('right_title', '')
            left_points = slide_data.get('left_points', [])
            right_points = slide_data.get('right_points', [])
            slide = create_comparison_slide(prs, title, left_title, right_title,
                                          left_points, right_points, slide_color_scheme)

            # Add logo if auto_logos is enabled (comparison slides don't have backgrounds)
            if auto_logos and should_add_logo(slide_type, title):
                logo_path = select_logo_by_content(title, left_points + right_points)
                if logo_path:
                    add_logo_to_slide(slide, logo_path, 'bottom_right')

        elif slide_type == 'two_column':
            title = slide_data.get('title', '')
            left_content = slide_data.get('left_content', [])
            right_content = slide_data.get('right_content', [])
            slide = create_two_column_slide(prs, title, left_content, right_content, slide_color_scheme)

            # Add logo if auto_logos is enabled
            if auto_logos and should_add_logo(slide_type, title):
                logo_path = select_logo_by_content(title, left_content + right_content)
                if logo_path:
                    add_logo_to_slide(slide, logo_path, 'bottom_right')

        elif slide_type == 'blank':
            slide = create_blank_slide(prs)

            # Add text if specified
            if 'text' in slide_data:
                add_text_to_slide(slide, slide_data['text'],
                                left=slide_data.get('text_left', 1),
                                top=slide_data.get('text_top', 1),
                                font_size=slide_data.get('font_size', 18))

            # Add image if specified
            if 'image' in slide_data:
                image_path = slide_data['image']
                shadow = slide_data.get('image_shadow', True)
                add_image_with_effects(slide, image_path,
                                     left=slide_data.get('image_left', 1),
                                     top=slide_data.get('image_top', 1),
                                     width=slide_data.get('image_width'),
                                     height=slide_data.get('image_height'),
                                     shadow=shadow)

            # Add logo if auto_logos is enabled
            if auto_logos and not slide_data.get('image'):
                title_for_logo = slide_data.get('text', 'Blank Slide')
                logo_path = select_logo_by_content(title_for_logo)
                if logo_path:
                    add_logo_to_slide(slide, logo_path, 'bottom_right')

    prs.save(output_path)
    print(f"Successfully created presentation: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


def create_simple_presentation(title: str, output_path: str, slides: Optional[List[Dict]] = None,
                              color_scheme: str = 'feedmob', auto_backgrounds: bool = True,
                              auto_logos: bool = True):
    """Create a simple presentation with title slide and optional content slides."""
    prs = Presentation()

    # Apply color scheme
    apply_color_scheme(prs, color_scheme)

    # Create title slide with auto-selected background
    background_image = None
    if auto_backgrounds:
        background_image = select_background_by_content(title)

    create_title_slide(prs, title, color_scheme=color_scheme, background_image=background_image)

    # Add content slides if provided
    if slides:
        for slide_info in slides:
            slide_title = slide_info.get('title', '')
            slide_content = slide_info.get('content', [])
            if isinstance(slide_content, str):
                slide_content = [slide_content]

            # Auto-select background for content slides
            slide_background = None
            if auto_backgrounds and should_add_background('content', slide_title):
                slide_background = select_background_by_content(slide_title, slide_content)

            if slide_background:
                slide = create_content_slide_with_background(prs, slide_title, slide_content, color_scheme, slide_background)
            else:
                slide = create_content_slide(prs, slide_title, slide_content, color_scheme)

            # Add logo if no background and auto_logos is enabled
            if auto_logos and should_add_logo('content', slide_title, bool(slide_background)):
                logo_path = select_logo_by_content(slide_title, slide_content)
                if logo_path:
                    add_logo_to_slide(slide, logo_path, 'bottom_right')

    prs.save(output_path)
    print(f"Successfully created presentation: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    if auto_backgrounds:
        print("✓ Automatic background selection enabled")
    if auto_logos:
        print("✓ Automatic logo insertion enabled")


def main():
    parser = argparse.ArgumentParser(
        description='Create PowerPoint presentations using python-pptx with enhanced styling'
    )
    parser.add_argument(
        '--output', '-o',
        required=True,
        help='Output file path for the PPTX file'
    )
    parser.add_argument(
        '--title',
        help='Title for the presentation (creates title slide)'
    )
    parser.add_argument(
        '--json', '-j',
        help='JSON file containing slide definitions'
    )
    parser.add_argument(
        '--subtitle',
        help='Subtitle for the title slide'
    )
    parser.add_argument(
        '--slides',
        nargs='+',
        help='Additional slide titles (creates content slides)'
    )
    parser.add_argument(
        '--color-scheme',
        choices=['professional', 'modern', 'corporate', 'binance', 'feedmob'],
        default='feedmob',
        help='Color scheme for the presentation (default: feedmob)'
    )
    parser.add_argument(
        '--background-image',
        help='Background image for title slide'
    )
    parser.add_argument(
        '--no-auto-backgrounds',
        action='store_true',
        help='Disable automatic background selection based on content'
    )
    parser.add_argument(
        '--no-auto-logos',
        action='store_true',
        help='Disable automatic logo insertion for slides without backgrounds'
    )

    args = parser.parse_args()

    # Ensure output path has .pptx extension
    output_path = args.output
    if not output_path.endswith('.pptx'):
        output_path += '.pptx'

    try:
        # Set auto-backgrounds and auto-logos based on command line flags
        auto_backgrounds = not args.no_auto_backgrounds
        auto_logos = not args.no_auto_logos

        if args.json:
            # Create from JSON file
            with open(args.json, 'r', encoding='utf-8') as f:
                json_data = json.load(f)

            # Override auto settings if specified via command line
            json_data['auto_backgrounds'] = auto_backgrounds
            json_data['auto_logos'] = auto_logos

            create_presentation_from_json(json_data, output_path)

        elif args.title:
            # Create simple presentation
            slides_data = []
            if args.slides:
                for slide_title in args.slides:
                    slides_data.append({'title': slide_title, 'content': []})

            create_simple_presentation(args.title, output_path, slides_data,
                                     args.color_scheme, auto_backgrounds, auto_logos)

        else:
            print("Error: Either --title or --json must be provided")
            parser.print_help()
            sys.exit(1)

    except FileNotFoundError as e:
        print(f"Error: File not found: {e}")
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON format: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"Error creating presentation: {e}")
        sys.exit(1)


if __name__ == '__main__':
    main()
